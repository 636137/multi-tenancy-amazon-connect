#!/usr/bin/env python3
"""
Create a professional PowerPoint presentation summarizing the Amazon Connect project.
Uses python-pptx with AWS-inspired design: navy + orange color scheme.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Colors: AWS-inspired palette
NAVY = RGBColor(35, 47, 62)          # Deep navy (main background)
ORANGE = RGBColor(255, 153, 0)       # AWS orange (accent)
LIGHT_GRAY = RGBColor(242, 242, 242) # Off-white (cards/content)
WHITE = RGBColor(255, 255, 255)      # White (text on navy)
DARK_GRAY = RGBColor(51, 51, 51)     # Dark gray (body text)
TEAL = RGBColor(0, 153, 179)         # Teal (secondary accent)

def create_presentation():
    """Create and return a presentation object."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs

def add_title_slide(prs, title, subtitle):
    """Add a dark title slide with white text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.LEFT
    
    # Orange accent bar at bottom
    shape = slide.shapes.add_shape(1, Inches(0), Inches(5.3), Inches(10), Inches(0.325))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.color.rgb = ORANGE
    
    return slide

def add_content_slide(prs, title, content_items=None, use_dark_bg=False):
    """Add a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    bg_color = NAVY if use_dark_bg else LIGHT_GRAY
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE if use_dark_bg else NAVY
    
    # Orange accent line under title
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(2), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.color.rgb = ORANGE
    
    # Content
    if content_items:
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(3.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE if use_dark_bg else DARK_GRAY
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            p.level = 0
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Add a slide with two columns of content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    # Orange accent line
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(2), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.color.rgb = ORANGE
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(4))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(4))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)
    
    return slide

def add_large_stat_slide(prs, statistic, description):
    """Add a slide with a large statistic in the center."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    
    # Large stat
    stat_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    stat_frame = stat_box.text_frame
    stat_frame.word_wrap = True
    p = stat_frame.paragraphs[0]
    p.text = statistic
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1.5))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    p = desc_frame.paragraphs[0]
    p.text = description
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    """Create the complete presentation."""
    prs = create_presentation()
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "Amazon Connect\nVoice Testing & Agent Guides",
        "IRS Taxpayer Services Platform"
    )
    
    # Slide 2: Overview
    add_content_slide(
        prs,
        "What This Does",
        [
            "🔊 Real Phone Calls — Actual PSTN calls via AWS Chime SDK",
            "🤖 AI-Powered Caller — Lambda with Polly TTS + intelligent responses",
            "💬 AI-to-AI Conversations — Nova Sonic bidirectional streaming",
            "📋 Agent Step-by-Step Guides — Interactive Cards on call accept",
            "📊 Status Tracking — Real-time call state via DynamoDB"
        ],
        use_dark_bg=True
    )
    
    # Slide 3: Key Features
    add_two_column_slide(
        prs,
        "Key Capabilities",
        [
            "REAL PHONE CALLS",
            "✓ PSTN integration",
            "✓ Test IVR flows",
            "✓ Validate contact flows",
            "",
            "AGENT GUIDES",
            "✓ Interactive cards",
            "✓ 5 IRS use cases",
            "✓ Step-by-step procedures"
        ],
        [
            "AI POWERED",
            "✓ Intelligent responses",
            "✓ Natural conversation",
            "✓ Multi-turn dialogs",
            "",
            "AUTOMATION",
            "✓ Regression testing",
            "✓ Load testing",
            "✓ Deployment ready"
        ]
    )
    
    # Slide 4: IRS Use Cases
    add_content_slide(
        prs,
        "IRS Agent Guide Cards (5 Topics)",
        [
            "💰 Refund Status — E-file timeline, paper filing, IDRS codes",
            "🛡️ Identity Verification — Letter 5071C/6331C, ID.me options",
            "💳 Payment Plans — Installment agreements, Form 433",
            "📄 Transcripts — Form 4506-T, online portal, 3-year availability",
            "📬 Notices — CP2000, CP14, response deadlines"
        ],
        use_dark_bg=True
    )
    
    # Slide 5: How It Works
    add_large_stat_slide(
        prs,
        "📞 1 Call",
        "From Dialing to Agent Card Display in Real-Time"
    )
    
    # Slide 6: Technical Architecture
    add_content_slide(
        prs,
        "Technical Stack",
        [
            "☁️  Amazon Connect — Contact center platform",
            "🤖 AWS Lex — Conversational AI (V2)",
            "⚡ AWS Lambda — Serverless handlers & AI logic",
            "📞 AWS Chime SDK — Phone integration & WebRTC",
            "🧠 Amazon Bedrock — Next-gen AI models (Nova Sonic)",
            "📊 DynamoDB — Call state & customer tracking",
            "🔊 Amazon Polly — Text-to-speech for AI caller"
        ]
    )
    
    # Slide 7: Live Demo
    add_content_slide(
        prs,
        "Live Demo Instance",
        [
            "🌐 Treasury Connect (Production)",
            "",
            "📞 Phone: +1 833-289-6602",
            "👤 Agent Portal: [Connect Agent Workspace]",
            "",
            "✅ Direct-to-agent routing (no IVR)",
            "✅ IRS Agent Guide Cards appear on call accept",
            "✅ 5 interactive procedure cards ready to deploy"
        ],
        use_dark_bg=True
    )
    
    # Slide 8: Project Structure
    add_two_column_slide(
        prs,
        "Project Structure",
        [
            "CORE COMPONENTS",
            "/lambda — Handlers",
            "/lex — Lex bot definitions",
            "/contact_flows — Flow exports",
            "/voice_tester — Test framework",
            "",
            "SUPPORTING",
            "/cdk — Infrastructure as Code",
            "/scripts — Deployment utils",
            "/docs — Documentation"
        ],
        [
            "SPECIALIZED MODULES",
            "/irs-agent-views",
            "  - 5 use case cards",
            "  - Dashboard views",
            "  - Setup guides",
            "",
            "/census-ai-agent",
            "  - Survey agent",
            "  - Deepgram TTS/STT",
            "  - CloudFront dashboard"
        ]
    )
    
    # Slide 9: Getting Started
    add_content_slide(
        prs,
        "Getting Started",
        [
            "1️⃣  Run ./install.sh — Install dependencies",
            "2️⃣  Configure AWS credentials in .env",
            "3️⃣  Run ./deploy.sh — Deploy infrastructure (CDK)",
            "4️⃣  Test with voice_tester — Make real test calls",
            "5️⃣  Monitor in AWS Console — DynamoDB call logs"
        ]
    )
    
    # Slide 10: Key Documents
    add_content_slide(
        prs,
        "Documentation",
        [
            "📖 README.md — Complete overview & architecture",
            "🚀 QUICKSTART.md — 5-minute setup guide",
            "🔧 irs-agent-views/SETUP_GUIDE.md — Card configuration",
            "📝 DEPLOYMENT_SUCCESS.md — Post-deployment checklist",
            "🎯 Test Results Available — voice_output/ directory"
        ],
        use_dark_bg=True
    )
    
    # Slide 11: Closing
    add_title_slide(
        prs,
        "Ready to Deploy?",
        "Automated Voice Testing for Amazon Connect Contact Centers"
    )
    
    # Save presentation
    output_path = "/Users/ChadDHendren/AmazonConnect1/Project_Summary.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created: {output_path}")
    return output_path

if __name__ == "__main__":
    main()
